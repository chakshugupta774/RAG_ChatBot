# file: document_loader.py
import io
import pandas as pd
import docx
import pdfplumber
from pptx import Presentation
from typing import List, Dict

def load_document(file_name: str, file_bytes: bytes) -> List[Dict]:
    """Load document and return list of dicts with text and metadata"""
    ext = file_name.split(".")[-1].lower()
    if ext == "pdf":
        return extract_pdf(file_name, file_bytes)
    elif ext == "docx":
        return extract_docx(file_name, file_bytes)
    elif ext == "txt":
        return extract_txt(file_name, file_bytes)
    elif ext in ["xlsx", "csv"]:
        return extract_excel_csv(file_name, file_bytes, ext)
    elif ext == "pptx":
        return extract_pptx(file_name, file_bytes)
    else:
        return [{"text": "Unsupported file format.", "source": file_name, "type": ext}]

def extract_pdf(file_name, file_bytes):
    docs = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                docs.append({"text": text, "source": file_name, "type": "pdf", "page": i+1})
    return docs

def extract_docx(file_name, file_bytes):
    doc = docx.Document(io.BytesIO(file_bytes))
    text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    return [{"text": text, "source": file_name, "type": "docx"}] if text else []

def extract_txt(file_name, file_bytes):
    text = file_bytes.decode("utf-8")
    return [{"text": text, "source": file_name, "type": "txt"}] if text.strip() else []

def extract_excel_csv(file_name, file_bytes, ext):
    if ext == "csv":
        df = pd.read_csv(io.BytesIO(file_bytes))
    else:
        df = pd.read_excel(io.BytesIO(file_bytes), engine="openpyxl")
    text = df.to_string()
    return [{"text": text, "source": file_name, "type": ext}] if not df.empty else []

def extract_pptx(file_name, file_bytes):
    docs = []
    prs = Presentation(io.BytesIO(file_bytes))
    for i, slide in enumerate(prs.slides):
        slide_text = "\n".join([s.text for s in slide.shapes if hasattr(s, "text") and s.text.strip()])
        if slide_text:
            docs.append({"text": slide_text, "source": file_name, "type": "pptx", "slide": i+1})
    return docs
